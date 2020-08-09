from django.shortcuts import render, redirect
from django.http import HttpResponse


#For language detection
from langdetect import detect


from django.core.files.storage import FileSystemStorage
import os
from django.db import models
from django.core.files.storage import default_storage


from gtts import gTTS
#import operator
import urllib


#NlP Packages summarization
import nltk 
from nltk.corpus import stopwords 
from nltk.tokenize import word_tokenize, sent_tokenize

#Web Scrapping Pkg
from bs4 import BeautifulSoup



#for Text pre-processing
import re
import unicodedata


#For internet connection checking
import requests
import urllib.request as urllib2
from urllib.request import urlopen

from django.conf.urls.static import static
from django.conf import settings


#Fir File Handeling
import PyPDF2
from pptx import Presentation
import glob



# Create your views here.

def index(request):
	return render(request, 'home.html')




def internet_on():
	try:
		urllib2.urlopen('http://216.58.192.142', timeout=60)
		return True
	except urllib2.URLError as err:
		return False



def ManualTxt(request):
	if request.method == 'POST':
		text=request.POST['fulltextarea']
		if detect(text)=='en' or len(text)>=50:

			if request.POST['submit'] == 'Convert_to_voice':
				return render(request,'ManualText.html',{'text':text})





			if request.POST['submit'] == 'Save':
				if internet_on():
					name=request.POST['File_name']
					if name == "":
						tts = gTTS(text)
						tts.save(f'media/audio/random.mp3')
						return render(request, 'home.html')
					else:
						
						tts = gTTS(text)
						tts.save(f'media/audio/{name}.mp3')
						return render(request, 'home.html')
				else:
					return render(request, "error.html")






			if request.POST['submit'] == 'Summarize_Text':
				if internet_on():
					stopWords = set(stopwords.words("english")) 
					words = word_tokenize(text) 
					   
					# Creating a frequency table to keep the  
					# score of each word 
					   
					freqTable = dict() 
					for word in words: 
					    word = word.lower() 
					    if word in stopWords: 
					        continue
					    if word in freqTable: 
					        freqTable[word] += 1
					    else: 
					        freqTable[word] = 1
					   
					# Creating a dictionary to keep the score 
					# of each sentence 
					sentences = sent_tokenize(text) 
					sentenceValue = dict() 
					   
					for sentence in sentences: 
					    for word, freq in freqTable.items(): 
					        if word in sentence.lower(): 
					            if sentence in sentenceValue: 
					                sentenceValue[sentence] += freq 
					            else: 
					                sentenceValue[sentence] = freq 
					   
					   
					   
					sumValues = 0
					for sentence in sentenceValue: 
					    sumValues += sentenceValue[sentence] 
					   
					# Average value of a sentence from the original text 
					   
					average = int(sumValues / len(sentenceValue)) 
					   
					# Storing sentences into our summary. 
					summary = '' 
					for sentence in sentences: 
					    if (sentence in sentenceValue) and (sentenceValue[sentence] > (1.2 * average)): 
					        summary += " " + sentence 
					return render(request, "home.html",{'text':summary})


				else:
					return render(request, "error.html")

		else:
			msg="Sorry! You can only enter the English text and length should be more than 50 characters"
			return render(request, "home.html",{'msg':msg, 'text':text})

		
			

			
				

					
#WEB SCRAPPING



def External_Text(request):
	if request.method == 'POST':
		if request.POST['submit'] == 'Scrap':
			raw_text=request.POST['fulltextarea']
			if internet_on():
				page=urlopen(raw_text)
				soup=BeautifulSoup(page,'lxml')
				fetched_text=' '.join(map(lambda p:p.text,soup.find_all('p')))

				if detect(fetched_text)=='en' or len(fetched_text)>=1000:
					fetched_text= unicodedata.normalize('NFKD', fetched_text).encode('ascii', 'ignore').decode('utf-8', 'ignore')
					def remove_special_characters(text, remove_digits=False):
						test = re.sub(r'[\[].*[\]]', '', text)
						punctuation = """!\"#$%&()*+-/;<=>?@^_`{|}~"""
						edited = ""
						for i in test:
							if i not in punctuation:
						 		edited += i
						return edited

					Actual_text= remove_special_characters(fetched_text, remove_digits=True)
					return render(request,'home.html', {'text': Actual_text})

				else:
					msg="Sorry! You can only enter the English text and length should be more than 1000 characters"
					return render(request, "home.html",{'msg':msg, 'text':text})

			else:
				return render(request, 'error.html')


					



				

			

def File_Upload(request):
	txxt=""
	if request.method=='POST':
		uploaded_file=request.FILES['document']
		name=uploaded_file.name
		nam=name[-5:]
		nm=name[-4:]
		url=f"media/ThisFile{nam}"
		url1=f"media/ThisFile{nm}"
		if os.path.isfile(url):
			os.remove(url)
		elif os.path.isfile(url1):
			os.remove(url1)
		if nam==".pptx" or nm==".pdf" or nm==".txt":

 	#For pptx file handeling
			if nam=='.pptx':
				txt=""
				fs=FileSystemStorage()
				name=f"ThisFile{nam}"
				fs.save(name, uploaded_file)
				url = fs.url(name)
				for eachfile in glob.glob(url):
				    prs = Presentation(eachfile)
				    for slide in prs.slides:
				        for shape in slide.shapes:
				            if hasattr(shape, "text"):
				                txxt +=shape.text
				
				
				return render(request, 'home.html', {'text':txxt})



		#For pdf file handeling


			elif nm=='.pdf':
				fs=FileSystemStorage()
				name=f"ThisFile{nm}"
				fs.save(name, uploaded_file)
				url = fs.url(name)
				pdfFileObject = open(url, 'rb')
				pdfReader = PyPDF2.PdfFileReader(pdfFileObject)
				count = pdfReader.numPages
				if count>=30:
					x=25
					count=30
				else:
					count=20
					x=10	
				for i in range(x,count):
				    page = pdfReader.getPage(i)
				    txxt +=page.extractText()

				
				return render(request, 'home.html', {'text':txxt})


			elif nm=='.txt':
				fs=FileSystemStorage()
				name=f"ThisFile{nm}"
				fs.save(name, uploaded_file)
				url = fs.url(name)
				with open(url) as myfile:
					data = myfile.read()
	
				return render(request, 'home.html',{'text':data})

			
		else:
			msg="Sorry! You can upload .pptx, .pdf or .txt files only:("
			return render(request, 'home.html', {'msg':msg})


	










